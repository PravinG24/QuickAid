from pptx import Presentation
from pathlib import Path
import shutil
from datetime import datetime

pptx_path = Path(r"D:\Documents\MyMahir Docs\Course Work\Project Submissions\QuickAid_Presentation .pptx.pptx")
if not pptx_path.exists():
    raise SystemExit(f"ERROR: PPTX not found: {pptx_path}")

backup = pptx_path.with_name(pptx_path.stem + '_backup_' + datetime.now().strftime('%Y%m%d_%H%M%S') + pptx_path.suffix)
shutil.copy2(pptx_path, backup)
print('Backup created:', backup)

prs = Presentation(str(pptx_path))

# Notes mapping by 1-based slide index
notes = {
    1: 'Welcome: Project name "QuickAid" and team; one-line summary: cloud-native campus helpdesk for students and staff.',
    2: 'Agenda: Problem → Objectives → Design → Implementation → Testing → Deployment → Demo → Limitations → Future work → Sources.',
    3: 'Problem statement: student/staff friction with manual support channels; slow responses, lost requests, no central tracking.',
    4: 'Objectives: provide online ticket submission & tracking; role-based access (User/Agent/Admin); automated notifications.',
    5: 'Research highlights: cloud suitability for education; serverless fits variable campus workloads.',
    6: 'Requirements: functional (submit/retrieve/update tickets, notifications, admin assignment) and non-functional (availability, low cost, security).',
    7: 'High-level architecture: Frontend (App Service), Backend (Azure Functions), DB (Cosmos DB), Blob store, Key Vault, notifications.',
    8: 'Technology stack: Azure App Service, Azure Functions, Cosmos DB, Blob Storage, Key Vault, MS Entra ID, SendGrid/Azure Comm Services, App Insights.',
    9: 'Backend design: event-driven functions for create/update, notifications, activity log persistence; partitioning rationale for Cosmos DB.',
    10: 'Frontend design: SPA flows for ticket form, status lookup, admin console; MSAL/Entra ID for admin, JWT/local for students.',
    11: 'Data model: ticket fields and partition key choice (/email); activity log stored per ticket.',
    12: 'Security & secrets: Key Vault for secrets; DefaultAzureCredential and role enforcement for Admin/Agent/User routes.',
    13: 'Notifications & activity logs: blob-backed notification store and SendGrid for email delivery; activity entries stamped and persisted.',
    14: 'Admin console: view tickets, bulk updates, assignment, approvals; Entra ID for admin auth.',
    15: 'User journeys: submit -> confirm email -> track; agent: view assigned -> update -> add notes; admin: assign -> bulk update.',
    16: 'Implementation milestones: core APIs, frontend, auth, notifications, deployment pipeline completed; backups and test artifacts available.',
    17: 'Limitations: no advanced reporting yet; some features depend on Azure quotas; external integrations not implemented.',
    18: 'Future work: reporting & exports, ML ticket classification, SLA dashboards, richer workflows.',
    19: 'Project timeline: Research→Design→Build→Test→Deploy→Present; weekly checkpoints and owners noted.',
    20: 'Demo setup: live Azure deployment used; sample accounts prepared for User/Agent/Admin.',
    21: 'Demo step 1: submit ticket through frontend; show confirmation email.',
    22: 'Demo step 2: track & view ticket details and activity log in real time.',
    23: 'Demo step 3: admin actions - assign ticket, update status to In Progress.',
    24: 'Test environment: testing ran in live Azure; tools: Azure API Management and live frontend; roles tested: User, Agent, Admin.',
    25: 'Testing approach & results: functional testing of core features; ticket submission, retrieval, status updates, emails, admin assignment; all core tests passed.',
    26: 'Issues & observations: TC-08 (assignment rejected due to category/team mismatch); TC-06 role enforcement blocked unauthorized updates; TC-12 prevented double-finishing; 22 test cases documented.',
    27: 'Monitoring & logging: Application Insights for telemetry, error tracking, request traces; example trace shows end-to-end ticket submission flow.',
    28: 'Cost & scalability: serverless reduces costs for bursty usage; plan Cosmos DB RU throughput for production.',
    29: 'Deployment: frontend on App Service, backend on Functions, Cosmos DB for data, Azure Comm Services/SendGrid for emails, Key Vault, App Insights; portal confirms resources.',
    30: 'Demonstration full flow: submit -> assign -> update -> resolve -> notification; show App Insights traces for validation.',
    31: 'User feedback & validation: summary of demo user feedback and acceptance criteria.',
    32: 'Ethical & compliance: minimal PII storage, secure connections, follow university tenant policies for guest access.',
    33: 'Sources/References: compiled from provided task resources and literature; full list in the project write-up.',
    34: 'Acknowledgements: advisors, teammates, university tenant, Azure credits.',
    35: 'Q&A: invite questions and offer to repeat specific flows.'
}

for i, slide in enumerate(prs.slides, start=1):
    note_text = notes.get(i, '')
    notes_slide = slide.notes_slide
    if notes_slide is None:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            # skip if cannot get notes_slide
            continue
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        continue
    if note_text:
        # replace existing notes
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = note_text
        print(f'Updated notes for slide {i}')
    else:
        print(f'No note provided for slide {i}; leaving existing notes')

# save
try:
    prs.save(str(pptx_path))
    print('Saved notes into', pptx_path)
except PermissionError:
    alt = pptx_path.with_name(pptx_path.stem + '_with_notes' + pptx_path.suffix)
    prs.save(str(alt))
    print('Original file locked. Saved annotated copy to', alt)
    print('Backup is at:', backup)
except Exception as e:
    print('ERROR saving PPTX:', e)
    print('Backup is at:', backup)
    raise
